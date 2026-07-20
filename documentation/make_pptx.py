from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

KERBLUE  = RGBColor(0x43, 0x38, 0xCA)
KERCORAL = RGBColor(0xE0, 0x5C, 0x52)
KERGREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARKGRAY = RGBColor(0x1F, 0x2A, 0x44)
LIGHTBG  = RGBColor(0xF8, 0xF9, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

def add_rect(slide, l, t, w, h, fill=None, line=None, radius=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARKGRAY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def slide_bg(slide, color=LIGHTBG):
    bg = add_rect(slide, 0, 0, 13.33, 7.5, fill=color)
    return bg

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.4, fill=KERBLUE)
    add_text(slide, title, 0.4, 0.15, 12, 0.75, size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12, 0.45, size=16, color=RGBColor(0xCC,0xCC,0xFF), align=PP_ALIGN.LEFT)
    # logo simulé
    logo = add_rect(slide, 12.4, 0.25, 0.7, 0.7, fill=KERCORAL)
    add_text(slide, "KC", 12.4, 0.3, 0.7, 0.6, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def footer(slide, page_num, total=18):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=KERBLUE)
    add_text(slide, "KerConnect  —  Naratechvision  —  contact@naratechvision.com", 0.3, 7.12, 11, 0.3,
             size=10, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, f"{page_num} / {total}", 12.5, 7.12, 0.7, 0.3, size=10, color=WHITE, align=PP_ALIGN.RIGHT)

def bullet_box(slide, bullets, l, t, w, h, title=None, bg=None, title_color=KERBLUE):
    if bg:
        add_rect(slide, l, t, w, h, fill=bg)
    if title:
        add_text(slide, title, l+0.15, t+0.1, w-0.3, 0.45, size=15, bold=True, color=title_color)
        t_offset = 0.55
    else:
        t_offset = 0.1
    for i, b in enumerate(bullets):
        add_text(slide, f"•  {b}", l+0.15, t+t_offset+i*0.42, w-0.3, 0.4, size=12, color=DARKGRAY)

def icon_card(slide, emoji, label, desc, l, t, w=2.5, h=2.2, bg=LIGHTBG, border=KERBLUE):
    add_rect(slide, l, t, w, h, fill=bg, line=border)
    add_text(slide, emoji, l, t+0.1, w, 0.6, size=28, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t+0.7, w, 0.45, size=13, bold=True, color=KERBLUE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, l+0.1, t+1.15, w-0.2, 0.9, size=10, color=DARKGRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 1 — TITRE
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, fill=KERBLUE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=KERCORAL)
add_rect(sl, 0, 7.42, 13.33, 0.08, fill=KERCORAL)
# Cercles décoratifs
add_rect(sl, 9.5, -0.5, 4, 4, fill=RGBColor(0x50,0x45,0xD5))
add_rect(sl, -0.5, 4, 3, 3, fill=RGBColor(0x50,0x45,0xD5))
# Logo KC
logo_bg = add_rect(sl, 5.67, 0.8, 2, 2, fill=WHITE)
add_text(sl, "KC", 5.67, 0.9, 2, 1.8, size=60, bold=True, color=KERCORAL, align=PP_ALIGN.CENTER)
add_text(sl, "KerConnect", 0.5, 3.2, 12.33, 1.1, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Plateforme Immobilière Sénégalaise", 0.5, 4.35, 12.33, 0.7, size=24, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
add_text(sl, "Naratechvision  •  contact@naratechvision.com  •  Juin 2026", 0.5, 5.5, 12.33, 0.5, size=14, color=RGBColor(0xAA,0xBB,0xFF), align=PP_ALIGN.CENTER)
add_rect(sl, 4, 5.1, 5.33, 0.05, fill=KERCORAL)

# ─────────────────────────────────────────────
# SLIDE 2 — SOMMAIRE
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Sommaire")
footer(sl, 2)
items = [
    ("01", "Le Problème", "Le marché immobilier sénégalais"),
    ("02", "La Solution", "KerConnect — concept et valeur"),
    ("03", "Les Profils", "4 types d'utilisateurs"),
    ("04", "Fonctionnalités", "Tour de la plateforme"),
    ("05", "Flux Contrat", "De la demande au paiement"),
    ("06", "Modèle Économique", "Commission + revenus"),
    ("07", "Stack Technique", "Technologies utilisées"),
    ("08", "Roadmap", "Phases de développement"),
]
cols = [0.4, 6.9]
for i, (num, title, desc) in enumerate(items):
    col = cols[i % 2]
    row = 1.7 + (i // 2) * 1.3
    add_rect(sl, col, row, 6, 1.1, fill=WHITE, line=KERBLUE)
    add_rect(sl, col, row, 0.6, 1.1, fill=KERBLUE)
    add_text(sl, num, col, row+0.25, 0.6, 0.6, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, col+0.7, row+0.05, 4.8, 0.5, size=16, bold=True, color=KERBLUE)
    add_text(sl, desc,  col+0.7, row+0.55, 4.8, 0.45, size=12, color=DARKGRAY)

# ─────────────────────────────────────────────
# SLIDE 3 — LE PROBLÈME
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Le Problème", "Le marché immobilier sénégalais souffre de...")
footer(sl, 3)
problems = [
    ("❌", "Manque de\ntransparence", "Offres invérifiées, prix\nnon standardisés"),
    ("📄", "Faux contrats\net litiges", "Pas de traçabilité,\naccords informels"),
    ("🔍", "Accès difficile\naux offres", "Bouche-à-oreille,\npas de catalogue central"),
    ("💸", "Paiements non\ntracés", "Espèces sans reçu,\nimpossible à contester"),
]
for i, (emoji, label, desc) in enumerate(problems):
    icon_card(sl, emoji, label, desc, 0.5 + i*3.1, 1.6, w=2.9, h=2.5,
              bg=WHITE if i%2==0 else RGBColor(0xFF,0xF0,0xF0), border=KERCORAL)
add_rect(sl, 0.5, 4.4, 12.33, 0.06, fill=KERCORAL)
add_text(sl, "→  KerConnect digitalise et sécurise l'ensemble du parcours immobilier",
         0.5, 4.6, 12.33, 0.6, size=18, bold=True, color=KERBLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 4 — LA SOLUTION
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "La Solution KerConnect", "Une plateforme 100% en ligne, du bien au paiement")
footer(sl, 4)
solutions = [
    ("✅", "Annonces vérifiées", "Validation admin avant publication"),
    ("📝", "Contrat numérique", "Upload PDF + signature manuscrite + validation"),
    ("💳", "Paiements tracés", "Wave, Orange Money, espèces avec reçu PDF"),
    ("⭐", "Score de confiance", "Notes 1-5 étoiles + badge Vérifié"),
    ("🔔", "Alertes automatiques", "Notification sur critères personnalisés"),
    ("💰", "Commission auto", "5% prélevé sur chaque paiement confirmé"),
]
for i, (emoji, label, desc) in enumerate(solutions):
    col = 0.5 if i < 3 else 6.9
    row = 1.6 + (i % 3) * 1.5
    add_rect(sl, col, row, 6, 1.25, fill=WHITE, line=KERGREEN)
    add_text(sl, emoji, col+0.1, row+0.3, 0.7, 0.6, size=24, align=PP_ALIGN.CENTER)
    add_text(sl, label, col+0.9, row+0.1, 4.8, 0.5, size=14, bold=True, color=KERBLUE)
    add_text(sl, desc,  col+0.9, row+0.6, 4.8, 0.5, size=11, color=DARKGRAY)

# ─────────────────────────────────────────────
# SLIDE 5 — LES 4 PROFILS
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Les 4 Profils Utilisateurs")
footer(sl, 5)
profiles = [
    ("🏠", "Client", KERBLUE,
     ["Recherche biens à louer/acheter","Fait des demandes en ligne","Signe les contrats PDF","Paie via Wave/OM/espèces","Note le bailleur après le contrat"]),
    ("🏢", "Bailleur", KERGREEN,
     ["Publie des annonces (3 étapes)","Gère les demandes reçues","Envoie le contrat PDF","Confirme les paiements","Consulte son rapport mensuel"]),
    ("🏗️", "Propriétaire", RGBColor(0xD9,0x77,0x06),
     ["Vend des biens immobiliers","Mêmes capacités que bailleur","Interface partagée","Contrats de vente","Suivi des acheteurs"]),
    ("⚙️", "Admin (Hawoly)", RGBColor(0x7C,0x3A,0xED),
     ["Valide les annonces","Gère les utilisateurs","Configure la commission","Consulte les transactions","Voit les revenus plateforme"]),
]
for i, (emoji, label, color, bullets) in enumerate(profiles):
    l = 0.3 + i * 3.2
    add_rect(sl, l, 1.5, 3.0, 5.4, fill=WHITE, line=color)
    add_rect(sl, l, 1.5, 3.0, 0.8, fill=color)
    add_text(sl, emoji, l, 1.55, 3.0, 0.7, size=28, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(sl, label, l, 2.1, 3.0, 0.5, size=16, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    for j, b in enumerate(bullets):
        add_text(sl, f"• {b}", l+0.15, 2.4+j*0.57, 2.75, 0.5, size=10, color=DARKGRAY)

# ─────────────────────────────────────────────
# SLIDE 6 — FONCTIONNALITÉS PRINCIPALES
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Fonctionnalités Principales")
footer(sl, 6)
feats = [
    ("📢", "Publication annonce","3 étapes + upload photos","Validation admin"),
    ("🔍", "Catalogue biens","Filtres ville/budget/type","Location & Vente"),
    ("📋", "Gestion demandes","Accepter/Refuser","Messaging intégré"),
    ("📄", "Contrat PDF","Upload + Signature","Validation bailleur"),
    ("💳", "Paiement multi-mode","Wave/OM/Carte/Espèces","Reçu PDF auto"),
    ("📊", "Rapport mensuel","Revenus par mois","Export PDF"),
    ("⭐", "Score confiance","Notes 1-5 étoiles","Badge Vérifié admin"),
    ("🔔", "Alertes & Notifs","Polling 30 secondes","Email automatique"),
]
for i, (emoji, label, d1, d2) in enumerate(feats):
    col = (i % 4) * 3.2 + 0.3
    row = 1.6 + (i // 4) * 2.7
    add_rect(sl, col, row, 3.0, 2.3, fill=WHITE, line=KERBLUE)
    add_text(sl, emoji, col, row+0.05, 3.0, 0.75, size=30, align=PP_ALIGN.CENTER)
    add_text(sl, label, col, row+0.75, 3.0, 0.5, size=13, bold=True, color=KERBLUE, align=PP_ALIGN.CENTER)
    add_text(sl, d1, col, row+1.2, 3.0, 0.4, size=10, color=DARKGRAY, align=PP_ALIGN.CENTER)
    add_text(sl, d2, col, row+1.6, 3.0, 0.4, size=10, color=KERCORAL, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────
# SLIDE 7 — FLUX CONTRAT
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Flux Contrat — De la Demande au Paiement")
footer(sl, 7)
steps = [
    (KERBLUE,  "🙋", "Client fait\nune demande"),
    (KERGREEN, "✅", "Bailleur\naccepte"),
    (KERCORAL, "📄", "Bailleur envoie\nle PDF contrat"),
    (KERBLUE,  "✍️", "Client signe et\nrenvoie le PDF"),
    (KERGREEN, "👁️", "Bailleur valide\nla signature"),
    (KERCORAL, "💳", "Client paie\n(Wave/Espèces)"),
    (KERGREEN, "💰", "Commission\nprélevée 5%"),
]
for i, (color, emoji, label) in enumerate(steps):
    l = 0.5 + i * 1.8
    add_rect(sl, l, 1.7, 1.6, 1.6, fill=color)
    add_text(sl, emoji, l, 1.75, 1.6, 0.75, size=28, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(sl, label, l, 2.45, 1.6, 0.75, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, str(i+1), l+0.6, 1.7, 0.4, 0.35, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        add_text(sl, "→", l+1.6, 2.2, 0.2, 0.5, size=20, bold=True, color=KERBLUE, align=PP_ALIGN.CENTER)
# Détails en bas
details = [
    "Demande\nsoumise",
    "Contrat auto\ncréé",
    "Email envoyé\nau client",
    "Reçu après\nupload",
    "Statut:\nvalide",
    "Reçu PDF\nautomatique",
    "Net bailleur\nenregistré",
]
for i, d in enumerate(details):
    l = 0.5 + i * 1.8
    add_rect(sl, l+0.05, 3.5, 1.5, 0.8, fill=WHITE, line=RGBColor(0xCC,0xCC,0xFF))
    add_text(sl, d, l+0.05, 3.55, 1.5, 0.75, size=9, color=DARKGRAY, align=PP_ALIGN.CENTER)
# Légende
add_rect(sl, 0.3, 4.6, 12.5, 0.06, fill=KERCORAL)
add_text(sl, "Email automatique envoyé après chaque étape clé  •  Notifications temps réel dans la cloche",
         0.3, 4.8, 12.5, 0.5, size=13, color=KERBLUE, align=PP_ALIGN.CENTER, bold=True)

# ─────────────────────────────────────────────
# SLIDE 8 — MODES DE PAIEMENT
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Paiements Multi-Mode", "Sécurisés, tracés, avec reçu PDF automatique")
footer(sl, 8)
modes = [
    ("🌊", "Wave", "Paiement instantané\nMobile Money"),
    ("🟠", "Orange Money", "Réseau Orange\ntrès répandu au Sénégal"),
    ("💳", "Carte bancaire", "Visa / Mastercard\nvia PayDunya"),
    ("💵", "Espèces", "Remise en main propre\nConfirmation bailleur"),
    ("📄", "Chèque", "Chèque bancaire\nConfirmation bailleur"),
]
for i, (emoji, label, desc) in enumerate(modes):
    l = 0.5 + i * 2.5
    color = KERGREEN if i < 3 else KERCORAL
    add_rect(sl, l, 1.7, 2.3, 3.0, fill=WHITE, line=color)
    add_rect(sl, l, 1.7, 2.3, 0.7, fill=color)
    add_text(sl, emoji, l, 1.72, 2.3, 0.65, size=28, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(sl, label, l, 2.35, 2.3, 0.55, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, desc, l, 2.9, 2.3, 0.8, size=11, color=DARKGRAY, align=PP_ALIGN.CENTER)
    tag = "Auto (PayDunya)" if i < 3 else "Manuel (bailleur)"
    t_color = KERGREEN if i < 3 else KERCORAL
    add_text(sl, tag, l, 3.8, 2.3, 0.45, size=10, bold=True, color=t_color, align=PP_ALIGN.CENTER, italic=True)
add_rect(sl, 0.3, 5.0, 12.5, 0.06, fill=KERBLUE)
add_text(sl, "⚡ Blocage automatique si le contrat n'est pas encore validé par le bailleur",
         0.3, 5.2, 12.5, 0.5, size=14, bold=True, color=KERCORAL, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 9 — MODÈLE ÉCONOMIQUE
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Modèle Économique", "Commission automatique sur chaque transaction")
footer(sl, 9)
# Schéma commission
add_rect(sl, 0.5, 1.6, 5.5, 4.5, fill=WHITE, line=KERBLUE)
add_text(sl, "💰 Comment ça marche ?", 0.6, 1.7, 5.3, 0.6, size=16, bold=True, color=KERBLUE)
rows = [
    ("Loyer mensuel (brut)", "200 000 FCFA", DARKGRAY),
    ("Commission KerConnect (5%)", "−10 000 FCFA", KERCORAL),
    ("Net bailleur", "190 000 FCFA", KERGREEN),
]
for i, (label, val, color) in enumerate(rows):
    bg = RGBColor(0xF0,0xF4,0xFF) if i==2 else WHITE
    add_rect(sl, 0.55, 2.45+i*0.85, 5.4, 0.75, fill=bg, line=RGBColor(0xDD,0xDD,0xFF))
    add_text(sl, label, 0.65, 2.5+i*0.85, 3.5, 0.6, size=13, color=DARKGRAY)
    add_text(sl, val, 4.15, 2.5+i*0.85, 1.7, 0.6, size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)

# Projections
add_rect(sl, 6.3, 1.6, 6.5, 4.5, fill=WHITE, line=KERGREEN)
add_text(sl, "📈 Projections", 6.4, 1.7, 6.3, 0.6, size=16, bold=True, color=KERGREEN)
projs = [
    ("10 logements × 200K FCFA/mois", "→  100 000 FCFA/mois"),
    ("50 logements × 200K FCFA/mois", "→  500 000 FCFA/mois"),
    ("100 logements × 200K FCFA/mois","→  1 000 000 FCFA/mois"),
    ("Taux configurable : 1% à 50%", "Admin peut ajuster"),
    ("Commission désactivable", "Par simple toggle"),
]
for i, (label, val) in enumerate(projs):
    add_text(sl, label, 6.4, 2.45+i*0.8, 4.5, 0.4, size=11, color=DARKGRAY)
    add_text(sl, val,   6.4, 2.8+i*0.8, 6.0, 0.4, size=12, bold=True, color=KERGREEN)
add_rect(sl, 0.3, 6.3, 12.5, 0.06, fill=KERBLUE)
add_text(sl, "Revenus supplémentaires prévus : abonnements premium, annonces sponsorisées",
         0.3, 6.45, 12.5, 0.4, size=12, color=KERBLUE, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────
# SLIDE 10 — FONCTIONNALITÉS INNOVANTES
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Fonctionnalités Innovantes", "Ce qui différencie KerConnect de la concurrence")
footer(sl, 10)
innov = [
    ("⭐", "Score de Confiance", KERBLUE,
     "Notes 1-5 étoiles après chaque contrat finalisé.\nBadge 'Bailleur Vérifié' attribué par l'admin.\nNote moyenne visible sur chaque annonce."),
    ("🔔", "Alertes Biens", KERCORAL,
     "Client s'abonne à des critères (ville, budget, type).\nDès qu'un bien correspond → notification automatique.\nGestion depuis /client/alertes."),
    ("📊", "Rapport Mensuel", KERGREEN,
     "Bailleur génère son rapport en 1 clic.\nTotal brut, commission, net reçu par mois.\nExportable en PDF imprimable."),
    ("✉️", "Messagerie Intégrée", RGBColor(0x7C,0x3A,0xED),
     "Messages bailleur ↔ client depuis la plateforme.\nEnvoyés par email avec reply-to configuré.\nHistorique des échanges par demande."),
    ("📄", "Contrat PDF Numérique", KERCORAL,
     "Bailleur uploade son modèle PDF.\nClient télécharge, signe manuscritement, renvoie.\nBailleur valide avant d'autoriser le paiement."),
    ("💰", "Commission Automatique", KERGREEN,
     "Configurable par l'admin (défaut 5%).\nCalculée sur chaque paiement confirmé.\nTableau de bord des revenus plateforme."),
]
for i, (emoji, label, color, desc) in enumerate(innov):
    col = 0.3 if i < 3 else 6.8
    row = 1.55 + (i % 3) * 1.7
    add_rect(sl, col, row, 6.2, 1.5, fill=WHITE, line=color)
    add_rect(sl, col, row, 0.7, 1.5, fill=color)
    add_text(sl, emoji, col, row+0.35, 0.7, 0.75, size=26, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(sl, label, col+0.8, row+0.08, 5.2, 0.5, size=13, bold=True, color=color)
    add_text(sl, desc,  col+0.8, row+0.55, 5.2, 0.85, size=10, color=DARKGRAY)

# ─────────────────────────────────────────────
# SLIDE 11 — ARCHITECTURE TECHNIQUE
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Architecture Technique")
footer(sl, 11)
# Frontend box
add_rect(sl, 0.4, 1.6, 5.8, 4.5, fill=WHITE, line=KERBLUE)
add_rect(sl, 0.4, 1.6, 5.8, 0.6, fill=KERBLUE)
add_text(sl, "⚛️  Frontend", 0.5, 1.65, 5.6, 0.5, size=16, bold=True, color=WHITE)
fe_items = [
    ("Next.js 16.2.9", "Turbopack, App Router"),
    ("React 19", "Interface utilisateur"),
    ("TypeScript 5", "Typage statique"),
    ("Tailwind CSS 4", "Design système"),
    ("Zustand 5", "State (auth + cookies)"),
    ("TanStack Query v5", "Cache requêtes API"),
    ("proxy.ts", "Middleware Next.js 16"),
]
for i, (tech, desc) in enumerate(fe_items):
    add_text(sl, f"• {tech}", 0.55, 2.3+i*0.51, 2.5, 0.45, size=11, bold=True, color=KERBLUE)
    add_text(sl, desc, 3.1, 2.3+i*0.51, 2.9, 0.45, size=10, color=DARKGRAY)
# Backend box
add_rect(sl, 6.5, 1.6, 6.4, 4.5, fill=WHITE, line=KERCORAL)
add_rect(sl, 6.5, 1.6, 6.4, 0.6, fill=KERCORAL)
add_text(sl, "🐘  Backend + BDD", 6.6, 1.65, 6.2, 0.5, size=16, bold=True, color=WHITE)
be_items = [
    ("Laravel 13.8", "Framework PHP"),
    ("PHP 8.3", "Runtime"),
    ("Sanctum 4.3", "Tokens Bearer 7j"),
    ("MySQL 8.0+", "Base de données"),
    ("Laravel Storage", "Fichiers PDF/IMG"),
    ("Gmail SMTP TLS", "Emails OTP/notifs"),
    ("DomPDF", "Génération reçus"),
]
for i, (tech, desc) in enumerate(be_items):
    add_text(sl, f"• {tech}", 6.65, 2.3+i*0.51, 3.0, 0.45, size=11, bold=True, color=KERCORAL)
    add_text(sl, desc, 9.7, 2.3+i*0.51, 3.0, 0.45, size=10, color=DARKGRAY)
# Flèche entre les deux
add_text(sl, "⟷\nHTTP/JSON\nBearer Token", 5.6, 3.2, 1.0, 1.2, size=11, bold=True,
         color=KERBLUE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 12 — BASE DE DONNÉES
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Base de Données — 10 Tables Clés")
footer(sl, 12)
tables = [
    ("users", "id, name, email, role,\nnote_moyenne, verifie", KERBLUE),
    ("biens", "id, user_id, titre, nature,\nstatut, images(JSON)", KERBLUE),
    ("demandes", "id, bien_id, demandeur_id,\nstatut, prenom_nom", KERBLUE),
    ("contrats", "id, fichier_contrat,\nfichier_signe_client, statut", KERCORAL),
    ("paiements", "id, montant, commission,\nmontant_net, mode", KERCORAL),
    ("avis", "id, contrat_id, note 1-5,\nauteur_id, cible_id", KERGREEN),
    ("alertes", "id, user_id, ville,\nnature, prix_max", KERGREEN),
    ("settings", "cle, valeur (commission_taux,\ncommission_active)", KERGREEN),
    ("favoris", "user_id, bien_id", RGBColor(0x7C,0x3A,0xED)),
    ("otps", "email, code, expires_at,\nused, type", RGBColor(0x7C,0x3A,0xED)),
]
for i, (name, cols, color) in enumerate(tables):
    col = 0.3 + (i % 5) * 2.58
    row = 1.6 + (i // 5) * 2.7
    add_rect(sl, col, row, 2.45, 2.4, fill=WHITE, line=color)
    add_rect(sl, col, row, 2.45, 0.55, fill=color)
    add_text(sl, name, col, row+0.08, 2.45, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, cols, col+0.1, row+0.65, 2.25, 1.65, size=9, color=DARKGRAY)

# ─────────────────────────────────────────────
# SLIDE 13 — SÉCURITÉ
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Sécurité et Fiabilité")
footer(sl, 13)
sec_items = [
    ("🔐", "Tokens Sanctum", "Bearer Token TTL 7 jours\nRévocation à la déconnexion\nIntercepteur 401 auto"),
    ("📧", "OTP Email", "Code 4 chiffres par email\nExpiration 15 minutes\n3 tentatives max"),
    ("🛡️", "Middleware Roles", "role:client / bailleur / admin\nProtection toutes routes\nProxy.ts cookies"),
    ("🚦", "Rate Limiting", "5 tentatives login\nBlocage 60 secondes\n3 tentatives OTP"),
    ("📁", "Uploads Sécurisés", "Validation MIME stricte\nMax 20 Mo configuré\nStockage privé"),
    ("💳", "Paiement Bloqué", "Contrat.statut = valide requis\nVérification backend strict\nDouble protection"),
]
for i, (emoji, label, desc) in enumerate(sec_items):
    col = 0.4 + (i % 3) * 4.2
    row = 1.6 + (i // 3) * 2.5
    color = KERBLUE if i % 3 == 0 else (KERGREEN if i % 3 == 1 else KERCORAL)
    add_rect(sl, col, row, 3.9, 2.1, fill=WHITE, line=color)
    add_text(sl, emoji, col, row+0.1, 3.9, 0.7, size=28, align=PP_ALIGN.CENTER)
    add_text(sl, label, col, row+0.75, 3.9, 0.5, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, desc,  col+0.1, row+1.2, 3.7, 0.8, size=10, color=DARKGRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 14 — ROADMAP
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Roadmap — Phases de Développement")
footer(sl, 14)
phases = [
    ("V1", "Actuel", KERBLUE, [
        "✅ Inscription + OTP email",
        "✅ Catalogue biens vérifiés",
        "✅ Demandes + messagerie",
        "✅ Contrats PDF numériques",
        "✅ Paiements multi-mode",
        "✅ Commission automatique",
        "✅ Score confiance + alertes",
        "✅ Rapport mensuel bailleur",
    ]),
    ("V2", "Court terme", KERGREEN, [
        "📍 Visite virtuelle 3D",
        "📍 Comparateur de biens",
        "📍 Google Maps dynamique",
        "📍 Application mobile",
        "📍 Cautions en escrow",
        "📍 Abonnements premium",
    ]),
    ("V3", "Moyen terme", KERCORAL, [
        "🔮 IA recommandation biens",
        "🔮 PayDunya production",
        "🔮 Partenariats bancaires",
        "🔮 Syndication agences",
    ]),
    ("V4", "Long terme", RGBColor(0x7C,0x3A,0xED), [
        "🌍 Côte d'Ivoire",
        "🌍 Mali / Sénégal",
        "🌍 Cameroun",
        "🌍 Expansion régionale",
    ]),
]
for i, (ver, label, color, items) in enumerate(phases):
    l = 0.3 + i * 3.2
    add_rect(sl, l, 1.55, 3.0, 5.5, fill=WHITE, line=color)
    add_rect(sl, l, 1.55, 3.0, 0.9, fill=color)
    add_text(sl, ver, l, 1.58, 3.0, 0.55, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, label, l, 2.05, 3.0, 0.4, size=12, color=RGBColor(0xDD,0xEE,0xFF) if color==KERBLUE else WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(sl, item, l+0.1, 2.6+j*0.55, 2.8, 0.5, size=10, color=DARKGRAY)

# ─────────────────────────────────────────────
# SLIDE 15 — POURQUOI KERCONNECT
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Pourquoi KerConnect ?", "Les avantages compétitifs clés")
footer(sl, 15)
add_rect(sl, 0.3, 1.6, 12.5, 4.9, fill=WHITE, line=KERBLUE)
advantages = [
    ("🎯", "Spécialisé Sénégal", "Conçu pour le marché local : FCFA, Wave, Orange Money, villes sénégalaises"),
    ("🔒", "Sécurisé de bout en bout", "De la demande à la validation du contrat PDF signé, tout est tracé"),
    ("⚡", "Rapide à déployer", "Stack moderne (Next.js 16 + Laravel 13), prêt pour la production"),
    ("💡", "Modèle économique solide", "Commission automatique dès le 1er paiement, sans intervention manuelle"),
    ("📈", "Scalable", "Architecture API REST, peut gérer des milliers d'annonces et de paiements"),
    ("🤝", "Confiance", "Score de confiance, badge Vérifié, contrats officiels — pas de faux annonces"),
]
for i, (emoji, label, desc) in enumerate(advantages):
    col = 0.5 if i < 3 else 6.7
    row = 1.75 + (i % 3) * 1.5
    add_text(sl, emoji, col, row+0.25, 0.7, 0.7, size=26, align=PP_ALIGN.CENTER)
    add_text(sl, label, col+0.8, row+0.05, 5.5, 0.5, size=14, bold=True, color=KERBLUE)
    add_text(sl, desc,  col+0.8, row+0.5, 5.5, 0.85, size=11, color=DARKGRAY)

# ─────────────────────────────────────────────
# SLIDE 16 — CHIFFRES CLÉS
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Chiffres Clés")
footer(sl, 16)
stats = [
    ("4", "Profils utilisateurs", KERBLUE),
    ("10", "Tables de données", KERCORAL),
    ("40+", "Endpoints API", KERBLUE),
    ("5%", "Commission défaut", KERGREEN),
    ("7j", "Durée token", KERCORAL),
    ("30s", "Polling notifs", KERGREEN),
]
for i, (val, label, color) in enumerate(stats):
    col = 0.5 + (i % 3) * 4.1
    row = 1.6 + (i // 3) * 2.7
    add_rect(sl, col, row, 3.8, 2.4, fill=color)
    add_text(sl, val, col, row+0.2, 3.8, 1.3, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, label, col, row+1.5, 3.8, 0.7, size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 17 — CONTACTS ET DÉMO
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "Contact et Informations", "Prêt pour une démonstration")
footer(sl, 17)
add_rect(sl, 0.4, 1.6, 8, 5.0, fill=WHITE, line=KERBLUE)
add_text(sl, "📍 Naratechvision", 0.6, 1.75, 7.5, 0.7, size=20, bold=True, color=KERBLUE)
contacts = [
    ("✉️ Email", "contact@naratechvision.com"),
    ("🌐 Plateforme", "http://localhost:3000 (dev)"),
    ("👤 Admin démo", "demehawoly@gmail.com"),
    ("🔑 Mot de passe", "KerConnect2026!"),
    ("📄 Documentation", "kerconnect_documentation.pdf"),
]
for i, (label, val) in enumerate(contacts):
    add_text(sl, label, 0.7, 2.6+i*0.75, 2.5, 0.6, size=13, bold=True, color=KERCORAL)
    add_text(sl, val,   3.3, 2.6+i*0.75, 4.8, 0.6, size=13, color=DARKGRAY)
# QR Code zone
add_rect(sl, 8.8, 1.6, 4.0, 5.0, fill=KERBLUE)
add_text(sl, "🚀", 8.8, 2.0, 4.0, 1.5, size=60, align=PP_ALIGN.CENTER, color=WHITE)
add_text(sl, "Demander\nune démo", 8.8, 3.6, 4.0, 1.0, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "contact@naratechvision.com", 8.8, 4.8, 4.0, 0.6, size=12, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 18 — MERCI
# ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, fill=KERBLUE)
add_rect(sl, 0, 0, 13.33, 0.12, fill=KERCORAL)
add_rect(sl, 0, 7.38, 13.33, 0.12, fill=KERCORAL)
add_rect(sl, 4.5, 4.5, 4.33, 0.08, fill=KERCORAL)
add_text(sl, "Merci de votre attention !", 0.5, 1.5, 12.33, 1.5, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "KerConnect — Plateforme Immobilière Sénégalaise", 0.5, 3.1, 12.33, 0.8, size=22, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
add_text(sl, "Naratechvision  •  contact@naratechvision.com", 0.5, 4.7, 12.33, 0.6, size=16, color=RGBColor(0xAA,0xBB,0xFF), align=PP_ALIGN.CENTER)
add_text(sl, "Version 2.0 — Juin 2026", 0.5, 5.5, 12.33, 0.5, size=14, color=RGBColor(0x88,0x99,0xCC), align=PP_ALIGN.CENTER, italic=True)

# Sauvegarde
out = r"C:\Users\DELL INSPIRON 16\Desktop\Dossier important\Ker_connect\documentation\KerConnect_Presentation.pptx"
prs.save(out)
print(f"PPTX sauvegarde: {out}")
